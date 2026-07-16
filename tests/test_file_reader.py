from pathlib import Path
import zipfile

import pandas as pd
import benchcore.file_reader as file_reader
from pptx import Presentation

from benchcore.execution import ContainerRunner, RunResult
from benchcore.file_reader import (
    FileReadLimits,
    clear_file_reader_cache,
    file_reader_cache_info,
    read_file,
    read_file_result,
    search_file,
)


class TimeoutContainerRunner(ContainerRunner):
    def __init__(self):
        super().__init__("file-reader:test", engine="/bin/true")
        self.command = None
        self.policy = None
        self.workspace_files = []

    def run(self, command, policy=None):
        self.command = command
        self.policy = policy
        self.workspace_files = sorted(path.name for path in command.cwd.iterdir())
        return RunResult(
            argv=command.argv,
            exit_code=None,
            stdout="",
            stderr="",
            elapsed_seconds=policy.timeout_seconds,
            timed_out=True,
            isolation="ephemeral_container_readonly_workspace",
            backend="docker",
        )


class OutputContainerRunner(ContainerRunner):
    def __init__(self, stdout: str):
        super().__init__("file-reader:test", engine="/bin/true")
        self.stdout = stdout

    def run(self, command, policy=None):
        return RunResult(
            argv=command.argv,
            exit_code=0,
            stdout=self.stdout,
            stderr="",
            elapsed_seconds=0.01,
            timed_out=False,
            isolation="ephemeral_container_readonly_workspace",
            backend="docker",
        )


def test_plain_reader_keeps_head_and_tail_for_long_files(tmp_path: Path):
    path = tmp_path / "manual.txt"
    path.write_text(
        "important opening\n" + ("middle filler\n" * 300) + "emergency mutual aid agreement\n",
        encoding="utf-8",
    )

    text = read_file(path, 500)

    assert "important opening" in text
    assert "middle truncated" in text
    assert "emergency mutual aid agreement" in text


def test_xlsx_reader_and_search_cover_rows_beyond_preview_head(tmp_path: Path):
    path = tmp_path / "annual.xlsx"
    df = pd.DataFrame(
        {
            "company": [f"company-{i}" for i in range(200)] + ["山东凯马汽车"],
            "net_profit": [i for i in range(200)] + [-8343],
        }
    )
    df.to_excel(path, index=False)

    preview = read_file(path, 600)
    hits = search_file(path, ["山东凯马汽车", "-8343"])

    assert "middle truncated" in preview
    assert "山东凯马汽车" in preview
    assert hits["山东凯马汽车"] is not None
    assert hits["-8343"] is not None


def test_legacy_doc_reader_does_not_route_binary_doc_to_python_docx(tmp_path: Path):
    path = tmp_path / "legacy.doc"
    path.write_bytes(b"legacy")

    text = read_file(path)

    assert "legacy .doc parsing unavailable" in text
    assert "FILE_READER_STATUS=security_blocked" in text


def test_zip_bomb_compression_ratio_is_rejected_before_extraction(tmp_path: Path):
    path = tmp_path / "bomb.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "word/document.xml",
            "<w:document><w:t>" + ("A" * 1_000_000) + "</w:t></w:document>",
        )
    limits = FileReadLimits(
        max_file_bytes=2_000_000,
        max_total_uncompressed_bytes=2_000_000,
        max_member_uncompressed_bytes=2_000_000,
        max_compression_ratio=10,
    )

    result = read_file_result(path, limits=limits)

    assert result.status == "security_blocked"
    assert result.details["code"] == "archive_compression_ratio"


def test_archive_total_uncompressed_budget_is_enforced(tmp_path: Path):
    path = tmp_path / "oversize.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr("word/document.xml", "<w:t>ok</w:t>" + ("x" * 600))
        archive.writestr("word/header1.xml", "<w:t>head</w:t>" + ("y" * 600))
    limits = FileReadLimits(
        max_file_bytes=10_000,
        max_total_uncompressed_bytes=1_000,
        max_member_uncompressed_bytes=900,
    )

    result = read_file_result(path, limits=limits)

    assert result.status == "budget_exceeded"
    assert result.details["code"] == "total_uncompressed_budget"


def test_content_hash_cache_does_not_reuse_same_path_after_bytes_change(tmp_path: Path):
    clear_file_reader_cache()
    path = tmp_path / "same-name.txt"
    path.write_text("first content", encoding="utf-8")
    first = read_file_result(path)
    path.write_text("second bytes", encoding="utf-8")
    second = read_file_result(path)

    assert first.text == "first content"
    assert second.text == "second bytes"
    assert first.content_sha256 != second.content_sha256
    info = file_reader_cache_info()
    assert info.hits == 0
    assert info.misses == 2


def test_cache_reuses_identical_content_across_paths_but_limits_change_key(tmp_path: Path):
    clear_file_reader_cache()
    left = tmp_path / "left.txt"
    right = tmp_path / "right.txt"
    left.write_text("identical", encoding="utf-8")
    right.write_text("identical", encoding="utf-8")

    assert read_file_result(left).text == "identical"
    assert read_file_result(right).text == "identical"
    info = file_reader_cache_info()
    assert info.hits == 1 and info.misses == 1

    different_limits = FileReadLimits(max_extracted_chars=100)
    assert read_file_result(right, limits=different_limits).text == "identical"
    assert file_reader_cache_info().misses == 2


def test_file_and_output_budgets_are_explicit_not_silent(tmp_path: Path):
    path = tmp_path / "large.txt"
    path.write_text("head " + ("middle " * 200) + "tail", encoding="utf-8")

    rejected = read_file_result(path, limits=FileReadLimits(max_file_bytes=20))
    assert rejected.status == "budget_exceeded"
    assert rejected.details["code"] == "file_size_budget"

    partial_limits = FileReadLimits(max_extracted_chars=80)
    partial = read_file_result(path, limits=partial_limits)
    assert partial.status == "truncated"
    assert "middle truncated" in partial.text
    search = search_file(path, ["not-present"], limits=partial_limits)
    assert search["_status"] == "truncated"
    assert search["_code"] == "incomplete_search_space"


def test_external_converter_timeout_is_operational_and_workspace_is_minimal(tmp_path: Path):
    path = tmp_path / "legacy.doc"
    path.write_bytes(b"legacy")
    runner = TimeoutContainerRunner()

    result = read_file_result(path, runner=runner)

    assert result.status == "operational_failed"
    assert result.details["code"] == "external_parser_timeout"
    assert runner.workspace_files == ["input.doc"]
    assert runner.command.argv == ("antiword", "input.doc")
    assert str(path) not in " ".join(runner.command.argv)
    assert runner.policy.network_enabled is False
    assert runner.policy.max_output_chars <= FileReadLimits().max_extracted_chars + 4096


def test_ooxml_page_budget_is_explicitly_truncated(tmp_path: Path):
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    for label in ("first", "middle", "last"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = label
    presentation.save(path)

    result = read_file_result(path, limits=FileReadLimits(max_pages=2))

    assert result.status == "truncated"
    assert result.details["page_count"] == 3
    assert result.details["parsed_pages"] == 2
    assert "first" in result.text and "last" in result.text


def test_xlsx_shared_string_rich_text_runs_preserve_one_index_per_si(tmp_path: Path):
    path = tmp_path / "rich.xlsx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "xl/sharedStrings.xml",
            """<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
            <si><r><t>Rich</t></r><r><t xml:space="preserve"> Text</t></r></si>
            <si><t>Second</t></si></sst>""",
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            """<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
            <sheetData><row><c t="s"><v>0</v></c><c t="s"><v>1</v></c></row></sheetData>
            </worksheet>""",
        )

    result = read_file_result(path)

    assert result.status == "ok"
    assert "Rich Text\tSecond" in result.text
    assert "invalid shared-string index" not in result.text


def test_external_driver_metadata_makes_silent_truncation_fail_closed(tmp_path: Path):
    path = tmp_path / "report.pdf"
    path.write_bytes(b"not parsed by fake isolated runner")
    marker = (
        '__BENCHCORE_FILE_READER_META__'
        '{"parsed_count": 2, "source_count": 20, "total_chars": 9000, "truncated": true}'
    )
    runner = OutputContainerRunner(marker + "\n[page 1]\nvisible prefix")

    result = read_file_result(path, runner=runner)
    search = search_file(path, ["missing evidence"], runner=runner)

    assert result.status == "truncated"
    assert result.details["source_count"] == 20
    assert search["_status"] == "truncated"
    assert search["_code"] == "incomplete_search_space"

    clear_file_reader_cache()
    invalid = read_file_result(
        path,
        runner=OutputContainerRunner("[page 1]\ntext without protocol metadata"),
    )
    assert invalid.status == "operational_failed"
    assert invalid.details["code"] == "external_parser_protocol"
