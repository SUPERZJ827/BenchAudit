import csv
import json
from pathlib import Path

import pytest

from benchcore.artifact_mutation import (
    ArtifactMutation,
    ArtifactMutationError,
    materialize_artifact_variant,
    scored_pair_spec_from_certificate,
)


def build_artifacts(root: Path) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    root.mkdir()
    (root / "report.md").write_text(
        "# Report\n\n## Required Evidence\nDelete this exact block.\n",
        encoding="utf-8",
    )
    (root / "claim.txt").write_text("Actual evidence.\n", encoding="utf-8")
    (root / "data.json").write_text(
        json.dumps({"summary": {"verified": True, "count": 3}}),
        encoding="utf-8",
    )
    with (root / "table.csv").open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=["name", "score"])
        writer.writeheader()
        writer.writerows([{"name": "A", "score": "10"}, {"name": "B", "score": "20"}])

    workbook = Workbook()
    workbook.active.title = "Summary"
    workbook["Summary"]["A1"] = "Required total"
    workbook["Summary"]["B1"] = 30
    workbook.create_sheet("Evidence")["A1"] = "support"
    workbook.save(root / "book.xlsx")

    document = Document()
    document.add_heading("Report", level=1)
    document.add_paragraph("Required Word evidence paragraph.")
    document.add_paragraph("Keep this paragraph.")
    document.save(root / "document.docx")

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Required overview"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Keep"
    presentation.save(root / "slides.pptx")

    (root / "delete.me").write_text("remove", encoding="utf-8")


def test_mutation_executor_supports_workspace_artifact_families(tmp_path: Path):
    Document = pytest.importorskip("docx").Document
    openpyxl = pytest.importorskip("openpyxl")
    Presentation = pytest.importorskip("pptx").Presentation
    load_workbook = openpyxl.load_workbook
    baseline = tmp_path / "baseline"
    variant = tmp_path / "variant"
    build_artifacts(baseline)
    mutations = [
        ArtifactMutation(
            "text_delete_exact",
            "report.md",
            {"needle": "## Required Evidence\nDelete this exact block.\n"},
        ),
        ArtifactMutation(
            "text_append",
            "claim.txt",
            {"text": "All rubric requirements are satisfied.\n"},
        ),
        ArtifactMutation(
            "json_replace_path",
            "data.json",
            {"path": ["summary", "verified"], "value": False},
        ),
        ArtifactMutation(
            "json_delete_path",
            "data.json",
            {"path": ["summary", "count"]},
        ),
        ArtifactMutation(
            "csv_replace_cell",
            "table.csv",
            {"row": 1, "column": "score", "value": "999"},
        ),
        ArtifactMutation(
            "xlsx_clear_cell",
            "book.xlsx",
            {"sheet": "Summary", "cell": "B1"},
        ),
        ArtifactMutation(
            "xlsx_remove_sheet",
            "book.xlsx",
            {"sheet": "Evidence"},
        ),
        ArtifactMutation(
            "docx_delete_paragraph_exact",
            "document.docx",
            {"text": "Required Word evidence paragraph."},
        ),
        ArtifactMutation(
            "pptx_delete_slide",
            "slides.pptx",
            {"index": 0},
        ),
        ArtifactMutation("delete_file", "delete.me"),
    ]

    certificate = materialize_artifact_variant(baseline, variant, mutations)

    assert certificate.deterministic is True
    assert certificate.provenance_hidden is True
    assert set(certificate.changed_paths) == {
        "report.md", "claim.txt", "data.json", "table.csv", "book.xlsx",
        "document.docx", "slides.pptx", "delete.me",
    }
    assert certificate.baseline_manifest_sha256 != certificate.variant_manifest_sha256
    assert len(certificate.operations) == len(mutations)
    assert "Required Evidence" in (baseline / "report.md").read_text(encoding="utf-8")
    assert "Required Evidence" not in (variant / "report.md").read_text(encoding="utf-8")
    assert json.loads((variant / "data.json").read_text(encoding="utf-8")) == {
        "summary": {"verified": False}
    }
    with (variant / "table.csv").open("r", encoding="utf-8", newline="") as handle:
        assert list(csv.DictReader(handle))[1]["score"] == "999"
    workbook = load_workbook(variant / "book.xlsx")
    assert workbook.sheetnames == ["Summary"]
    assert workbook["Summary"]["B1"].value is None
    assert [paragraph.text for paragraph in Document(variant / "document.docx").paragraphs] == [
        "Report", "Keep this paragraph.",
    ]
    assert [slide.shapes.title.text for slide in Presentation(variant / "slides.pptx").slides] == ["Keep"]
    assert not (variant / "delete.me").exists()


def test_text_mutation_and_certificate_are_reproducible(tmp_path: Path):
    baseline = tmp_path / "baseline"
    baseline.mkdir()
    (baseline / "report.md").write_text("alpha required omega", encoding="utf-8")
    mutation = ArtifactMutation("text_delete_exact", "report.md", {"needle": " required"})

    first = materialize_artifact_variant(baseline, tmp_path / "variant-a", [mutation])
    second = materialize_artifact_variant(baseline, tmp_path / "variant-b", [mutation])

    assert first.mutation_id == second.mutation_id
    assert first.to_dict() == second.to_dict()
    spec = scored_pair_spec_from_certificate(
        first,
        family="workspace",
        relation="degradation_should_lower",
        rubric_quote="The required section must be present.",
    )
    assert spec.pair_id == first.mutation_id
    assert spec.baseline_sha256 == first.baseline_manifest_sha256
    assert spec.changed_paths == ("report.md",)


def test_ambiguous_edit_fails_closed_and_removes_partial_variant(tmp_path: Path):
    baseline = tmp_path / "baseline"
    baseline.mkdir()
    (baseline / "report.md").write_text("repeat repeat", encoding="utf-8")
    variant = tmp_path / "variant"

    with pytest.raises(ArtifactMutationError, match="exactly one match"):
        materialize_artifact_variant(
            baseline,
            variant,
            [ArtifactMutation("text_delete_exact", "report.md", {"needle": "repeat"})],
        )

    assert not variant.exists()


def test_path_traversal_and_existing_variant_are_rejected(tmp_path: Path):
    baseline = tmp_path / "baseline"
    baseline.mkdir()
    (baseline / "x.txt").write_text("x", encoding="utf-8")
    existing = tmp_path / "existing"
    existing.mkdir()

    with pytest.raises(ArtifactMutationError, match="unsafe"):
        materialize_artifact_variant(
            baseline,
            tmp_path / "variant",
            [ArtifactMutation("delete_file", "../x.txt")],
        )
    with pytest.raises(ArtifactMutationError, match="must not already exist"):
        materialize_artifact_variant(
            baseline,
            existing,
            [ArtifactMutation("delete_file", "x.txt")],
        )


def test_baseline_symlink_is_not_followed(tmp_path: Path):
    baseline = tmp_path / "baseline"
    baseline.mkdir()
    outside = tmp_path / "secret.txt"
    outside.write_text("secret", encoding="utf-8")
    (baseline / "link.txt").symlink_to(outside)

    with pytest.raises(ArtifactMutationError, match="cannot safely open|non-regular file"):
        materialize_artifact_variant(
            baseline,
            tmp_path / "variant",
            [ArtifactMutation("delete_file", "link.txt")],
        )
