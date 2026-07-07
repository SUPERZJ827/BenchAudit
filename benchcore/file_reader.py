"""Universal input-file reader for benchmark auditing.

If a task provides a file, the auditor must be able to read it. This dispatches
by file type to the right extractor (spreadsheets, Word, PowerPoint, PDF, text)
and returns a compact text profile the LLM auditors can reason over. Unknown
types fall back to a best-effort text read; the caller can then let an LLM decide
how to interpret them.
"""
from __future__ import annotations

import warnings
from functools import lru_cache
from pathlib import Path

warnings.filterwarnings("ignore")


def _xlsx(path: Path, max_chars: int) -> str:
    body = _xlsx_full_text(path)
    return _head_tail(body, max_chars)


def _xlsx_full_text(path: Path) -> str:
    import pandas as pd
    xl = pd.ExcelFile(path)
    out = [f"sheets={xl.sheet_names}"]
    for sh in xl.sheet_names:
        df = xl.parse(sh, header=None)
        out.append(f"-- sheet '{sh}' shape={df.shape} --")
        clean = df.fillna("").astype(str)
        for row in clean.itertuples(index=False, name=None):
            line = "\t".join(cell.strip() for cell in row if cell.strip())
            if line:
                out.append(line)
    return "\n".join(out)


def _docx(path: Path, max_chars: int) -> str:
    from docx import Document
    d = Document(str(path))
    paras = [p.text for p in d.paragraphs if p.text.strip()]
    tables = []
    for t in d.tables[:3]:
        for row in t.rows[:8]:
            tables.append(" | ".join(c.text for c in row.cells))
    body = "\n".join(paras[:60])
    if tables:
        body += "\n[tables]\n" + "\n".join(tables[:24])
    return f"({len(d.paragraphs)} paras, {len(d.tables)} tables)\n" + body[:max_chars]


def _pptx(path: Path, max_chars: int) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for i, slide in enumerate(prs.slides):
        s = " | ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip())
        if s:
            texts.append(f"[slide {i+1}] {s}")
    return f"({len(prs.slides)} slides)\n" + "\n".join(texts)[:max_chars]


def _pdf(path: Path, max_chars: int, max_pages: int = 8) -> str:
    import pdfplumber
    with pdfplumber.open(str(path)) as pdf:
        n = len(pdf.pages)
        head = list(range(min(n, max_pages)))
        tail = list(range(max(0, n - 2), n)) if n > max_pages else []
        indices = sorted(set(head + tail))
        txt = []
        for idx in indices:
            txt.append(f"[page {idx + 1}]\n{pdf.pages[idx].extract_text() or ''}")
        body = "\n".join(txt)
    truncated = n > len(indices)
    return (
        f"({n}-page PDF, preview_pages={[i + 1 for i in indices]}, truncated={str(truncated).lower()})\n"
        + _head_tail(body, max_chars)
    )


def _head_tail(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    if max_chars <= 200:
        return text[:max_chars]
    marker = f"\n...[middle truncated; total_chars={len(text)}]...\n"
    side = max(1, (max_chars - len(marker)) // 2)
    return text[:side] + marker + text[-side:]


def _plain(path: Path, max_chars: int) -> str:
    return _head_tail(path.read_text(encoding="utf-8", errors="replace"), max_chars)


_DISPATCH = {
    ".xlsx": _xlsx, ".xls": _xlsx,
    ".docx": _docx, ".doc": _docx,
    ".pptx": _pptx,
    ".pdf": _pdf,
    ".csv": _plain, ".txt": _plain, ".md": _plain, ".json": _plain,
}


def read_file(path, max_chars: int = 3000) -> str:
    """Return a compact text profile of any supported input file."""
    path = Path(path)
    if not path.exists():
        return f"[{path.name}] (文件不存在)"
    ext = path.suffix.lower()
    reader = _DISPATCH.get(ext)
    try:
        if reader is None:
            return f"FILE {path.name} (类型 {ext} 暂无专用解析器, 尝试纯文本)\n" + _plain(path, max_chars)
        return f"FILE {path.name} [{ext}]\n" + reader(path, max_chars)
    except Exception as e:
        return f"FILE {path.name} [{ext}] (读取失败: {e})"


@lru_cache(maxsize=256)
def _searchable_text(path_str: str, max_pages: int | None = None) -> str:
    path = Path(path_str)
    ext = path.suffix.lower()
    if ext == ".pdf":
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            pages = pdf.pages if max_pages is None else pdf.pages[:max_pages]
            return "\n".join((p.extract_text() or "") for p in pages)
    if ext in (".xlsx", ".xls"):
        return _xlsx_full_text(path)
    if ext in (".docx", ".doc"):
        from docx import Document
        return "\n".join(p.text for p in Document(str(path)).paragraphs)
    return read_file(path, 200000)


def search_file(path, terms: list[str], max_pages: int | None = None) -> dict:
    """Search a (large) file's full text for terms; returns term -> found context.

    Lets value-rubric verification locate specific figures in big documents
    (e.g. a 193-page annual report) without sending the whole file to an LLM.
    """
    path = Path(path)
    try:
        text = _searchable_text(str(path), max_pages)
    except Exception as e:
        return {"_error": str(e)}
    out = {}
    low = text.lower()
    for t in terms:
        idx = low.find(t.lower())  # case-insensitive: 'purchase order' must match 'Purchase Order'
        out[t] = text[max(0, idx - 40): idx + 60].replace("\n", " ") if idx >= 0 else None
    return out
