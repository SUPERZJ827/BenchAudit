"""Deterministic, provenance-safe mutations for rubric-scored artifacts.

The caller (or an LLM planner) may select an operator and target, but this
module performs the actual edit, verifies containment, and emits a sidecar
certificate binding the complete before/after workspace manifests.  Mutation
metadata is never written into the artifact shown to the evaluator.

Only explicit, exactly addressable edits are supported.  Free-form rewriting
belongs to a review-only LLM path and is intentionally absent here.
"""
from __future__ import annotations

import csv
import hashlib
import json
import os
import shutil
import stat
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence


SUPPORTED_OPERATORS = frozenset({
    "delete_file",
    "text_delete_exact",
    "text_append",
    "json_delete_path",
    "json_replace_path",
    "csv_replace_cell",
    "xlsx_clear_cell",
    "xlsx_remove_sheet",
    "docx_delete_paragraph_exact",
    "pptx_delete_slide",
})


class ArtifactMutationError(ValueError):
    """The requested mutation is unsafe, ambiguous, unsupported, or a no-op."""


@dataclass(frozen=True)
class ArtifactMutation:
    operator: str
    relative_path: str
    parameters: Mapping[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class MutationStepEvidence:
    operator: str
    relative_path: str
    before_sha256: str | None
    after_sha256: str | None
    details: Mapping[str, Any]


@dataclass(frozen=True)
class ArtifactMutationCertificate:
    schema_version: str
    mutation_id: str
    baseline_manifest_sha256: str
    variant_manifest_sha256: str
    changed_paths: tuple[str, ...]
    operations: tuple[MutationStepEvidence, ...]
    provenance_hidden: bool
    deterministic: bool
    baseline_file_count: int
    variant_file_count: int
    total_baseline_bytes: int
    total_variant_bytes: int

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass(frozen=True)
class _Manifest:
    files: Mapping[str, Mapping[str, Any]]
    sha256: str
    total_bytes: int


def scored_pair_spec_from_certificate(
    certificate: ArtifactMutationCertificate,
    *,
    family: str,
    relation: str,
    rubric_quote: str,
    target_criterion: str | None = None,
    expected_min_delta: float = 0.10,
    invariance_tolerance: float = 0.03,
    explicit_requirement: bool = True,
    official_evaluator: bool = True,
    grader_kind: str = "llm",
):
    """Bind a materialized artifact mutation to the shared scored-pair spec."""

    from .counterexample_validation import ScoredPairSpec

    operators = "+".join(step.operator for step in certificate.operations)
    return ScoredPairSpec(
        pair_id=certificate.mutation_id,
        family=family,
        relation=relation,
        mutation_operator=operators,
        construction="deterministic" if certificate.deterministic else "hybrid",
        baseline_sha256=certificate.baseline_manifest_sha256,
        variant_sha256=certificate.variant_manifest_sha256,
        changed_paths=certificate.changed_paths,
        rubric_quote=rubric_quote,
        target_criterion=target_criterion,
        expected_min_delta=expected_min_delta,
        invariance_tolerance=invariance_tolerance,
        explicit_requirement=explicit_requirement,
        provenance_hidden=certificate.provenance_hidden,
        official_evaluator=official_evaluator,
        grader_kind=grader_kind,
    )


def materialize_artifact_variant(
    baseline_root: Path,
    variant_root: Path,
    mutations: Sequence[ArtifactMutation],
    *,
    max_files: int = 2_000,
    max_total_bytes: int = 512_000_000,
) -> ArtifactMutationCertificate:
    """Copy a baseline workspace, apply controlled edits, and certify the delta.

    ``variant_root`` must not exist and may not be inside ``baseline_root``.
    Baseline symlinks, devices, sockets, and FIFOs are rejected rather than
    followed.  On any failure, only the newly created variant directory is
    removed.
    """

    baseline = baseline_root.expanduser().resolve(strict=True)
    variant = variant_root.expanduser().absolute()
    if not baseline.is_dir():
        raise ArtifactMutationError("baseline_root must be a directory")
    if variant.exists() or variant.is_symlink():
        raise ArtifactMutationError("variant_root must not already exist")
    variant_resolved = variant.resolve(strict=False)
    if variant_resolved == baseline or variant_resolved.is_relative_to(baseline):
        raise ArtifactMutationError("variant_root may not be inside baseline_root")
    if not mutations:
        raise ArtifactMutationError("at least one mutation is required")
    if len(mutations) > 32:
        raise ArtifactMutationError("at most 32 mutation steps are allowed")
    _validate_mutation_requests(mutations)

    before = _workspace_manifest(baseline, max_files=max_files, max_total_bytes=max_total_bytes)
    variant.mkdir(parents=True, exist_ok=False)
    try:
        _copy_manifest_files(baseline, variant, before)
        steps = tuple(_apply_mutation(variant, mutation) for mutation in mutations)
        after = _workspace_manifest(variant, max_files=max_files, max_total_bytes=max_total_bytes)
        changed = tuple(sorted(
            path
            for path in set(before.files) | set(after.files)
            if before.files.get(path) != after.files.get(path)
        ))
        if not changed:
            raise ArtifactMutationError("mutation sequence did not change the workspace")
        expected_paths = {mutation.relative_path for mutation in mutations}
        unexpected = set(changed) - expected_paths
        if unexpected:
            raise ArtifactMutationError(
                "mutation changed undeclared path(s): " + ", ".join(sorted(unexpected))
            )
        payload = {
            "schema_version": "artifact-mutation-certificate-v1",
            "baseline_manifest_sha256": before.sha256,
            "variant_manifest_sha256": after.sha256,
            "changed_paths": changed,
            "operations": [asdict(step) for step in steps],
        }
        mutation_id = _canonical_sha256(payload)[:24]
        return ArtifactMutationCertificate(
            schema_version="artifact-mutation-certificate-v1",
            mutation_id=mutation_id,
            baseline_manifest_sha256=before.sha256,
            variant_manifest_sha256=after.sha256,
            changed_paths=changed,
            operations=steps,
            provenance_hidden=True,
            deterministic=True,
            baseline_file_count=len(before.files),
            variant_file_count=len(after.files),
            total_baseline_bytes=before.total_bytes,
            total_variant_bytes=after.total_bytes,
        )
    except Exception:
        shutil.rmtree(variant)
        raise


def _validate_mutation_requests(mutations: Sequence[ArtifactMutation]) -> None:
    for index, mutation in enumerate(mutations):
        if mutation.operator not in SUPPORTED_OPERATORS:
            raise ArtifactMutationError(
                f"mutation {index} uses unsupported operator: {mutation.operator}"
            )
        path = Path(mutation.relative_path)
        if (
            not mutation.relative_path
            or path.is_absolute()
            or ".." in path.parts
            or path.as_posix() in {".", ""}
        ):
            raise ArtifactMutationError(f"mutation {index} relative_path is unsafe")
        if not isinstance(mutation.parameters, Mapping):
            raise ArtifactMutationError(f"mutation {index} parameters must be an object")


def _workspace_manifest(root: Path, *, max_files: int, max_total_bytes: int) -> _Manifest:
    rows: dict[str, dict[str, Any]] = {}
    total = 0
    for directory, dirnames, filenames in os.walk(root, followlinks=False):
        directory_path = Path(directory)
        for name in sorted(dirnames):
            path = directory_path / name
            if path.is_symlink():
                raise ArtifactMutationError(f"symlink directory is not allowed: {path.relative_to(root)}")
        for name in sorted(filenames):
            path = directory_path / name
            relative = path.relative_to(root).as_posix()
            file_sha256, file_bytes = _regular_file_digest(path, relative)
            total += file_bytes
            if len(rows) + 1 > max_files:
                raise ArtifactMutationError(f"workspace exceeds {max_files} files")
            if total > max_total_bytes:
                raise ArtifactMutationError(f"workspace exceeds {max_total_bytes} bytes")
            rows[relative] = {
                "sha256": file_sha256,
                "bytes": file_bytes,
            }
    payload = [
        {"path": path, **rows[path]}
        for path in sorted(rows)
    ]
    return _Manifest(rows, _canonical_sha256(payload), total)


def _copy_manifest_files(source: Path, target: Path, manifest: _Manifest) -> None:
    for relative in sorted(manifest.files):
        src = source / relative
        dst = target / relative
        dst.parent.mkdir(parents=True, exist_ok=True)
        _copy_regular_file(
            src,
            dst,
            relative,
            expected_sha256=str(manifest.files[relative]["sha256"]),
            expected_bytes=int(manifest.files[relative]["bytes"]),
        )


def _apply_mutation(root: Path, mutation: ArtifactMutation) -> MutationStepEvidence:
    path = _contained_path(root, mutation.relative_path)
    before = _file_sha256(path) if path.is_file() else None
    operator = mutation.operator
    parameters = dict(mutation.parameters)
    if operator != "delete_file" and not path.is_file():
        raise ArtifactMutationError(f"target file does not exist: {mutation.relative_path}")

    if operator == "delete_file":
        if not path.is_file() or path.is_symlink():
            raise ArtifactMutationError(f"delete_file target is not a regular file: {mutation.relative_path}")
        _require_only(parameters, set(), operator)
        path.unlink()
        details = {"deleted": True}
    elif operator == "text_delete_exact":
        _require_only(parameters, {"needle"}, operator)
        needle = _nonempty_string(parameters.get("needle"), "needle")
        text = path.read_text(encoding="utf-8")
        count = text.count(needle)
        if count != 1:
            raise ArtifactMutationError(f"text_delete_exact requires exactly one match, found {count}")
        path.write_text(text.replace(needle, "", 1), encoding="utf-8")
        details = {"deleted_text_sha256": _canonical_sha256(needle), "matched": 1}
    elif operator == "text_append":
        _require_only(parameters, {"text"}, operator)
        addition = _nonempty_string(parameters.get("text"), "text")
        with path.open("a", encoding="utf-8") as handle:
            handle.write(addition)
        details = {"appended_text_sha256": _canonical_sha256(addition), "characters": len(addition)}
    elif operator in {"json_delete_path", "json_replace_path"}:
        allowed = {"path"} if operator == "json_delete_path" else {"path", "value"}
        _require_only(parameters, allowed, operator)
        keys = _json_path(parameters.get("path"))
        value = json.loads(path.read_text(encoding="utf-8"))
        parent, key = _json_parent(value, keys)
        if operator == "json_delete_path":
            previous = _delete_json_value(parent, key)
        else:
            previous = _replace_json_value(parent, key, parameters.get("value"))
        path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        details = {
            "json_path": keys,
            "previous_value_sha256": _canonical_sha256(previous),
        }
    elif operator == "csv_replace_cell":
        _require_only(parameters, {"row", "column", "value"}, operator)
        row_index = _nonnegative_int(parameters.get("row"), "row")
        column = parameters.get("column")
        with path.open("r", encoding="utf-8", newline="") as handle:
            rows = list(csv.DictReader(handle))
            fieldnames = list(rows[0].keys()) if rows else []
        if row_index >= len(rows):
            raise ArtifactMutationError("CSV row is out of range")
        if not isinstance(column, str) or column not in fieldnames:
            raise ArtifactMutationError("CSV column is missing")
        previous = rows[row_index][column]
        replacement = str(parameters.get("value", ""))
        if previous == replacement:
            raise ArtifactMutationError("CSV replacement is a no-op")
        rows[row_index][column] = replacement
        with path.open("w", encoding="utf-8", newline="") as handle:
            writer = csv.DictWriter(handle, fieldnames=fieldnames)
            writer.writeheader()
            writer.writerows(rows)
        details = {"row": row_index, "column": column, "previous_sha256": _canonical_sha256(previous)}
    elif operator == "xlsx_clear_cell":
        _require_only(parameters, {"sheet", "cell"}, operator)
        try:
            from openpyxl import load_workbook
        except ImportError as exc:  # pragma: no cover - depends on installation
            raise ArtifactMutationError("xlsx operators require openpyxl") from exc

        sheet = _nonempty_string(parameters.get("sheet"), "sheet")
        cell = _nonempty_string(parameters.get("cell"), "cell")
        workbook = load_workbook(path)
        if sheet not in workbook.sheetnames:
            raise ArtifactMutationError(f"Excel sheet does not exist: {sheet}")
        target = workbook[sheet][cell]
        if target.value is None:
            raise ArtifactMutationError("Excel cell is already empty")
        previous = target.value
        target.value = None
        workbook.save(path)
        details = {"sheet": sheet, "cell": cell, "previous_sha256": _canonical_sha256(previous)}
    elif operator == "xlsx_remove_sheet":
        _require_only(parameters, {"sheet"}, operator)
        try:
            from openpyxl import load_workbook
        except ImportError as exc:  # pragma: no cover - depends on installation
            raise ArtifactMutationError("xlsx operators require openpyxl") from exc

        sheet = _nonempty_string(parameters.get("sheet"), "sheet")
        workbook = load_workbook(path)
        if sheet not in workbook.sheetnames:
            raise ArtifactMutationError(f"Excel sheet does not exist: {sheet}")
        if len(workbook.sheetnames) <= 1:
            raise ArtifactMutationError("cannot remove the only Excel sheet")
        workbook.remove(workbook[sheet])
        workbook.save(path)
        details = {"removed_sheet": sheet}
    elif operator == "docx_delete_paragraph_exact":
        _require_only(parameters, {"text"}, operator)
        try:
            from docx import Document
        except ImportError as exc:  # pragma: no cover - depends on installation
            raise ArtifactMutationError("docx operators require python-docx") from exc

        text = _nonempty_string(parameters.get("text"), "text")
        document = Document(path)
        matches = [paragraph for paragraph in document.paragraphs if paragraph.text == text]
        if len(matches) != 1:
            raise ArtifactMutationError(
                f"docx_delete_paragraph_exact requires exactly one match, found {len(matches)}"
            )
        element = matches[0]._element
        element.getparent().remove(element)
        document.save(path)
        details = {"deleted_paragraph_sha256": _canonical_sha256(text)}
    elif operator == "pptx_delete_slide":
        _require_only(parameters, {"index"}, operator)
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - depends on installation
            raise ArtifactMutationError("pptx operators require python-pptx") from exc

        index = _nonnegative_int(parameters.get("index"), "index")
        presentation = Presentation(path)
        if index >= len(presentation.slides):
            raise ArtifactMutationError("PowerPoint slide index is out of range")
        if len(presentation.slides) <= 1:
            raise ArtifactMutationError("cannot remove the only PowerPoint slide")
        slide = presentation.slides[index]
        title = slide.shapes.title.text if slide.shapes.title is not None else ""
        slide_id = presentation.slides._sldIdLst[index]
        presentation.part.drop_rel(slide_id.rId)
        del presentation.slides._sldIdLst[index]
        presentation.save(path)
        details = {"removed_slide_index": index, "title_sha256": _canonical_sha256(title)}
    else:  # Defensive: validation already checked the allow-list.
        raise ArtifactMutationError(f"unsupported operator: {operator}")

    after = _file_sha256(path) if path.is_file() else None
    if before == after:
        raise ArtifactMutationError(f"{operator} produced no byte-level change")
    return MutationStepEvidence(operator, mutation.relative_path, before, after, details)


def _contained_path(root: Path, relative: str) -> Path:
    path = root / relative
    resolved = path.resolve(strict=False)
    if resolved == root or not resolved.is_relative_to(root):
        raise ArtifactMutationError(f"target escapes variant root: {relative}")
    if path.is_symlink():
        raise ArtifactMutationError(f"symlink target is not allowed: {relative}")
    return path


def _json_path(value: Any) -> list[str | int]:
    if not isinstance(value, list) or not value:
        raise ArtifactMutationError("JSON path must be a non-empty list")
    result: list[str | int] = []
    for part in value:
        if isinstance(part, bool) or not isinstance(part, (str, int)):
            raise ArtifactMutationError("JSON path components must be strings or integers")
        result.append(part)
    return result


def _json_parent(value: Any, keys: Sequence[str | int]) -> tuple[Any, str | int]:
    current = value
    for key in keys[:-1]:
        if isinstance(key, int) and isinstance(current, list) and 0 <= key < len(current):
            current = current[key]
        elif isinstance(key, str) and isinstance(current, dict) and key in current:
            current = current[key]
        else:
            raise ArtifactMutationError(f"JSON path does not exist at {key!r}")
    return current, keys[-1]


def _delete_json_value(parent: Any, key: str | int) -> Any:
    if isinstance(key, int) and isinstance(parent, list) and 0 <= key < len(parent):
        return parent.pop(key)
    if isinstance(key, str) and isinstance(parent, dict) and key in parent:
        return parent.pop(key)
    raise ArtifactMutationError("JSON delete target does not exist")


def _replace_json_value(parent: Any, key: str | int, replacement: Any) -> Any:
    if isinstance(key, int) and isinstance(parent, list) and 0 <= key < len(parent):
        previous = parent[key]
        if previous == replacement:
            raise ArtifactMutationError("JSON replacement is a no-op")
        parent[key] = replacement
        return previous
    if isinstance(key, str) and isinstance(parent, dict) and key in parent:
        previous = parent[key]
        if previous == replacement:
            raise ArtifactMutationError("JSON replacement is a no-op")
        parent[key] = replacement
        return previous
    raise ArtifactMutationError("JSON replacement target does not exist")


def _require_only(parameters: Mapping[str, Any], allowed: set[str], operator: str) -> None:
    missing = allowed - set(parameters)
    unexpected = set(parameters) - allowed
    if missing:
        raise ArtifactMutationError(f"{operator} missing parameter(s): {', '.join(sorted(missing))}")
    if unexpected:
        raise ArtifactMutationError(f"{operator} unexpected parameter(s): {', '.join(sorted(unexpected))}")


def _nonempty_string(value: Any, name: str) -> str:
    if not isinstance(value, str) or not value:
        raise ArtifactMutationError(f"{name} must be a non-empty string")
    return value


def _nonnegative_int(value: Any, name: str) -> int:
    if isinstance(value, bool) or not isinstance(value, int) or value < 0:
        raise ArtifactMutationError(f"{name} must be a non-negative integer")
    return value


def _file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _regular_file_digest(path: Path, relative: str) -> tuple[str, int]:
    flags = os.O_RDONLY | getattr(os, "O_CLOEXEC", 0) | getattr(os, "O_NOFOLLOW", 0)
    try:
        descriptor = os.open(path, flags)
    except OSError as exc:
        raise ArtifactMutationError(f"cannot safely open regular file {relative}: {exc}") from exc
    digest = hashlib.sha256()
    total = 0
    try:
        metadata = os.fstat(descriptor)
        if not stat.S_ISREG(metadata.st_mode):
            raise ArtifactMutationError(f"non-regular file is not allowed: {relative}")
        while True:
            chunk = os.read(descriptor, 1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
            total += len(chunk)
    finally:
        os.close(descriptor)
    return digest.hexdigest(), total


def _copy_regular_file(
    source: Path,
    target: Path,
    relative: str,
    *,
    expected_sha256: str,
    expected_bytes: int,
) -> None:
    flags = os.O_RDONLY | getattr(os, "O_CLOEXEC", 0) | getattr(os, "O_NOFOLLOW", 0)
    try:
        source_fd = os.open(source, flags)
    except OSError as exc:
        raise ArtifactMutationError(f"cannot safely reopen {relative}: {exc}") from exc
    target_fd: int | None = None
    digest = hashlib.sha256()
    total = 0
    try:
        metadata = os.fstat(source_fd)
        if not stat.S_ISREG(metadata.st_mode):
            raise ArtifactMutationError(f"source changed to a non-regular file: {relative}")
        target_fd = os.open(
            target,
            os.O_WRONLY | os.O_CREAT | os.O_EXCL | getattr(os, "O_CLOEXEC", 0),
            0o600,
        )
        while True:
            chunk = os.read(source_fd, 1024 * 1024)
            if not chunk:
                break
            view = memoryview(chunk)
            while view:
                written = os.write(target_fd, view)
                if written <= 0:
                    raise ArtifactMutationError(f"short write while copying: {relative}")
                view = view[written:]
            digest.update(chunk)
            total += len(chunk)
    finally:
        os.close(source_fd)
        if target_fd is not None:
            os.close(target_fd)
    if total != expected_bytes or digest.hexdigest() != expected_sha256:
        raise ArtifactMutationError(f"source changed while copying: {relative}")


def _canonical_sha256(value: Any) -> str:
    payload = json.dumps(value, sort_keys=True, ensure_ascii=False, separators=(",", ":"), default=str)
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()
